"""
Gerador da apresentação PowerPoint — S07 Qualidade de Software
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ──────────────────────────────────────────────────────────────────────────────
# PALETA
# ──────────────────────────────────────────────────────────────────────────────
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # azul-escuro quase preto
C_BLUE      = RGBColor(0x16, 0x21, 0x3E)   # azul médio
C_ACCENT    = RGBColor(0x0F, 0x3C, 0x78)   # azul accent
C_GREEN     = RGBColor(0x05, 0xA6, 0x77)   # verde Cypress
C_ORANGE    = RGBColor(0xFF, 0x6C, 0x37)   # laranja Postman
C_YELLOW    = RGBColor(0xFF, 0xD7, 0x00)   # amarelo Pokemon
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xE8, 0xF0, 0xFE)
C_GRAY      = RGBColor(0xA0, 0xA8, 0xB8)
C_CODE_BG   = RGBColor(0x0D, 0x11, 0x17)
C_CODE_GRN  = RGBColor(0x6A, 0xD0, 0x8A)
C_CODE_BLU  = RGBColor(0x79, 0xC0, 0xFF)
C_CODE_YLW  = RGBColor(0xE3, 0xB3, 0x41)
C_CODE_CMT  = RGBColor(0x8B, 0x94, 0x9E)
C_RED       = RGBColor(0xE5, 0x53, 0x4B)
C_SLIDE_BG  = RGBColor(0x0E, 0x17, 0x2B)


# ──────────────────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]   # completamente em branco
    return prs.slides.add_slide(layout)


def bg(slide, color=C_SLIDE_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill_color, radius=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def multiline_txt(slide, lines, l, t, w, h, size=14, color=C_WHITE, bold=False, spacing=1.0):
    """lines = list of (text, bold, color, size)"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, bold, color, size)
        line_text, line_bold, line_color, line_size = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size)
        run.font.bold = line_bold
        run.font.color.rgb = line_color
    return tb


def divider(slide, t, color=C_ACCENT, l=0.5, w=12.33):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def pill(slide, text, l, t, w, h, bg_color, txt_color=C_WHITE, size=13, bold=True):
    r = rect(slide, l, t, w, h, bg_color)
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = txt_color
    return r


def section_badge(slide, label, color):
    """small colored label in top-right"""
    pill(slide, label, 11.0, 0.15, 2.0, 0.38, color, C_WHITE, 11, True)


def add_table(slide, headers, rows, l, t, w, h,
              hdr_bg=C_ACCENT, hdr_fg=C_WHITE,
              row_bg=C_BLUE, row_alt=C_DARK, row_fg=C_WHITE,
              hdr_size=13, row_size=12):
    from pptx.util import Inches, Pt
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    # header
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(hdr_size)
        run.font.bold = True
        run.font.color.rgb = hdr_fg
    # rows
    for ri, row in enumerate(rows):
        bg_c = row_bg if ri % 2 == 0 else row_alt
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_c
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(row_size)
            run.font.color.rgb = row_fg
    return tbl


def code_block(slide, lines, l, t, w, h):
    """lines = list of (text, color)"""
    r = rect(slide, l, t, w, h, C_CODE_BG)
    tb = slide.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.12),
                                   Inches(w - 0.3), Inches(h - 0.24))
    tb.word_wrap = False
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for text, color in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(11)
        run.font.color.rgb = color
        run.font.name = "Courier New"
    return r


def arrow_right(slide, l, t, color=C_GRAY):
    txt(slide, "→", l, t, 0.3, 0.3, size=16, color=color)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — CAPA
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s, C_DARK)

# gradient-like accent bar left
rect(s, 0, 0, 0.55, 7.5, C_ACCENT)
rect(s, 0.55, 0, 0.12, 7.5, RGBColor(0x12, 0x4A, 0x9A))

# title area
txt(s, "S07", 1.2, 1.3, 11, 1.3, size=96, bold=True, color=C_WHITE)
txt(s, "QUALIDADE DE SOFTWARE", 1.25, 2.65, 10, 0.7, size=28, bold=True, color=C_LIGHT)
divider(s, 3.5, C_ACCENT, 1.25, 10.5)

txt(s, "Automação de Testes na Prática", 1.25, 3.65, 10, 0.55,
    size=20, color=C_GRAY)

# two project pills
pill(s, "🌲  Cypress  E2E", 1.25, 4.5, 3.2, 0.6, C_GREEN, size=16)
pill(s, "📮  Postman  +  Newman", 4.75, 4.5, 4.0, 0.6, C_ORANGE, size=16)
pill(s, "⚡  GitHub  Actions  CI/CD", 9.05, 4.5, 3.5, 0.6, C_ACCENT, size=16)

txt(s, "Material didático — Testes automatizados, API testing e integração contínua",
    1.25, 6.5, 11, 0.5, size=13, color=C_GRAY, italic=True)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — ESTRUTURA DO REPOSITÓRIO
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_ACCENT)
txt(s, "Estrutura do Repositório", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "VISÃO GERAL", C_ACCENT)

# tree visual
tree_lines = [
    ("📁  S07/",              C_WHITE,   True,  15),
    ("",                       C_WHITE,   False, 12),
    ("├── 📄  README.md",      C_CODE_YLW, False, 13),
    ("│",                       C_GRAY,    False, 13),
    ("├── 🌲  cypress-project/",  C_GREEN,   True,  14),
    ("│   ├── package.json",   C_CODE_BLU, False, 12),
    ("│   ├── cypress.config.js", C_CODE_BLU, False, 12),
    ("│   ├── cypress.bdd.config.js", C_CODE_BLU, False, 12),
    ("│   └── cypress/",       C_GREEN,   False, 12),
    ("│       ├── pages/          ← Page Objects", C_CODE_CMT, False, 11),
    ("│       ├── sem-bdd/specs/  ← Testes POM",   C_CODE_CMT, False, 11),
    ("│       ├── bdd/features/   ← Gherkin",       C_CODE_CMT, False, 11),
    ("│       ├── bdd/step_definitions/",            C_CODE_CMT, False, 11),
    ("│       ├── support/commands.js",              C_CODE_CMT, False, 11),
    ("│       └── fixtures/dados.json",              C_CODE_CMT, False, 11),
    ("│",                       C_GRAY,    False, 13),
    ("├── 📮  api-testing/",    C_ORANGE,  True,  14),
    ("│   ├── package.json",    C_CODE_BLU, False, 12),
    ("│   ├── collections/pokeapi.collection.json", C_CODE_YLW, False, 12),
    ("│   └── environments/pokeapi.environment.json", C_CODE_YLW, False, 12),
    ("│",                       C_GRAY,    False, 13),
    ("└── ⚡  .github/workflows/ci.yml", C_YELLOW, True, 14),
]

tb = slide_tb = s.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5.8), Inches(6.3))
tb.word_wrap = False
tf = tb.text_frame
tf.word_wrap = False
first = True
for text, color, bold, size in tree_lines:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Courier New"

# legend cards on right
def legend_card(slide, l, t, title, desc, icon, color):
    rect(slide, l, t, 6.0, 1.1, color)
    txt(slide, icon + "  " + title, l + 0.2, t + 0.08, 5.6, 0.4,
        size=16, bold=True, color=C_WHITE)
    txt(slide, desc, l + 0.2, t + 0.5, 5.6, 0.5, size=12, color=C_LIGHT)

legend_card(s, 6.9, 0.9,  "cypress-project/",
            "Testes E2E — Page Object Model + BDD/Gherkin\nAlvo: globalsqa.com",
            "🌲", C_GREEN)
legend_card(s, 6.9, 2.2,  "api-testing/",
            "Testes de API — Postman Collection + Newman\nAlvo: pokeapi.co",
            "📮", C_ORANGE)
legend_card(s, 6.9, 3.5,  ".github/workflows/ci.yml",
            "Pipeline CI/CD — 3 jobs paralelos\nDisparado a cada push para main",
            "⚡", C_ACCENT)
legend_card(s, 6.9, 4.8,  "README.md",
            "Documentação raiz — guia de início rápido\npara ambos os projetos",
            "📄", RGBColor(0x4A, 0x55, 0x68))


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — CYPRESS: VISÃO GERAL
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_GREEN)
txt(s, "Projeto 1 — Cypress E2E", 0.4, 0.1, 9, 0.55, size=26, bold=True)
section_badge(s, "CYPRESS", C_GREEN)

txt(s, "Site testado:  globalsqa.com/demo-site/", 0.4, 0.82, 9, 0.4,
    size=14, color=C_GRAY, italic=True)

# two approaches
rect(s, 0.3, 1.3, 5.9, 5.7, C_BLUE)
rect(s, 6.7, 1.3, 6.0, 5.7, RGBColor(0x0D, 0x22, 0x42))

txt(s, "Abordagem 1", 0.55, 1.4, 5.4, 0.4, size=11, color=C_GRAY)
txt(s, "Sem BDD", 0.55, 1.75, 5.4, 0.5, size=22, bold=True, color=C_GREEN)
txt(s, "Page Object Model", 0.55, 2.2, 5.4, 0.4, size=14, color=C_LIGHT)

items1 = [
    "📄  specs/*.spec.js  →  arquivo de teste",
    "🏗️  pages/*.js  →  Page Object (seletores + ações)",
    "🔧  support/commands.js  →  comandos reutilizáveis",
    "▶️  npm run test:sem-bdd",
]
for i, item in enumerate(items1):
    txt(s, item, 0.55, 2.75 + i * 0.72, 5.4, 0.55, size=13, color=C_WHITE)

txt(s, "Abordagem 2", 6.95, 1.4, 5.5, 0.4, size=11, color=C_GRAY)
txt(s, "Com BDD", 6.95, 1.75, 5.5, 0.5, size=22, bold=True, color=C_CODE_YLW)
txt(s, "Gherkin + Cucumber", 6.95, 2.2, 5.5, 0.4, size=14, color=C_LIGHT)

items2 = [
    "📋  features/*.feature  →  cenário Gherkin",
    "🔗  step_definitions/*.steps.js  →  implementação",
    "🏗️  pages/*.js  →  mesmo Page Object!",
    "▶️  npm run test:bdd",
]
for i, item in enumerate(items2):
    txt(s, item, 6.95, 2.75 + i * 0.72, 5.5, 0.55, size=13, color=C_WHITE)

# shared label
rect(s, 4.5, 6.65, 4.33, 0.28, RGBColor(0x0A, 0x2A, 0x5E))
txt(s, "⬆  Os dois estilos compartilham os mesmos Page Objects  ⬆",
    4.5, 6.65, 4.33, 0.28, size=10, color=C_GRAY, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — PAGE OBJECT MODEL
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_GREEN)
txt(s, "Page Object Model (POM)", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "CYPRESS — POM", C_GREEN)

# diagram boxes
def box(slide, l, t, w, h, color, title, subs):
    rect(slide, l, t, w, h, color)
    txt(slide, title, l + 0.15, t + 0.1, w - 0.3, 0.35,
        size=13, bold=True, color=C_WHITE)
    for i, sub in enumerate(subs):
        txt(slide, "•  " + sub, l + 0.15, t + 0.48 + i * 0.35, w - 0.3, 0.32,
            size=11, color=C_LIGHT)

box(s, 0.3, 0.85, 3.5, 2.2, C_ACCENT, "📄  alertas.spec.js",
    ["Importa AlertasPage",
     "describe() / it() / context()",
     "beforeEach → acessar()",
     "Chama métodos do Page Object"])

txt(s, "usa", 3.85, 1.7, 0.5, 0.3, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s, "→", 3.95, 1.55, 0.4, 0.4, size=22, color=C_GRAY, align=PP_ALIGN.CENTER)

box(s, 4.35, 0.85, 4.2, 2.2, RGBColor(0x05, 0x5A, 0x3C), "🏗️  AlertasPage.js",
    ["get botaoTentar → seletor CSS",
     "acessar() → cy.visit()",
     "clicarAba(i) → .eq(i).click()",
     "verificarResultado() → .should()"])

txt(s, "usa", 8.6, 1.7, 0.5, 0.3, size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s, "→", 8.7, 1.55, 0.4, 0.4, size=22, color=C_GRAY, align=PP_ALIGN.CENTER)

box(s, 9.1, 0.85, 3.9, 2.2, RGBColor(0x3A, 0x1C, 0x71), "🔧  commands.js",
    ["cy.acessarPagina()",
     "cy.acessarIframe()",
     "cy.aceitarConfirmacao()",
     "cy.responderPrompt()"])

divider(s, 3.3, C_GREEN)

# code sample left
txt(s, "spec.js  — O teste", 0.3, 3.45, 6.2, 0.35, size=13, bold=True, color=C_GREEN)
code_block(s, [
    ("import AlertasPage from '../../pages/AlertasPage'", C_CODE_YLW),
    ("", C_WHITE),
    ("describe('Caixas de Alerta', () => {", C_WHITE),
    ("  beforeEach(() => {", C_WHITE),
    ("    AlertasPage.acessar()", C_CODE_GRN),
    ("  })", C_WHITE),
    ("", C_WHITE),
    ("  it('deve exibir OK ao aceitar', () => {", C_CODE_BLU),
    ("    AlertasPage.clicarAba(1)", C_CODE_GRN),
    ("    AlertasPage.clicarTentarEAceitar()", C_CODE_GRN),
    ("    AlertasPage.verificarResultadoConfirmacao('You pressed OK!')", C_CODE_GRN),
    ("  })", C_WHITE),
    ("})", C_WHITE),
], 0.3, 3.85, 6.2, 3.35)

# code sample right
txt(s, "AlertasPage.js  — O Page Object", 6.8, 3.45, 6.2, 0.35, size=13, bold=True, color=C_GREEN)
code_block(s, [
    ("class AlertasPage {", C_WHITE),
    ("  // Getter — seletor CSS", C_CODE_CMT),
    ("  get botaoTentar() {", C_CODE_BLU),
    ("    return cy.get('.tab-pane.active button')", C_CODE_GRN),
    ("  }", C_WHITE),
    ("", C_WHITE),
    ("  acessar() {", C_CODE_BLU),
    ("    cy.acessarPagina('/demo-site/alertbox/')", C_CODE_GRN),
    ("  }", C_WHITE),
    ("", C_WHITE),
    ("  clicarTentarEAceitar() {", C_CODE_BLU),
    ("    cy.on('window:confirm', () => true)", C_CODE_YLW),
    ("    this.botaoTentar.click()", C_CODE_GRN),
    ("  }", C_WHITE),
    ("}  export default new AlertasPage()", C_WHITE),
], 6.8, 3.85, 6.2, 3.35)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — BDD / GHERKIN
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_GREEN)
txt(s, "BDD — Behavior Driven Development", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "CYPRESS — BDD", C_GREEN)

# flow arrows
boxes_flow = [
    ("📋\n.feature\n(Gherkin)", 0.3,  1.0, C_CODE_YLW),
    ("🔗\nstep_definitions\n(.steps.js)",  4.5, 1.0, C_ACCENT),
    ("🏗️\nPage Object\n(AlertasPage.js)", 8.7, 1.0, RGBColor(0x05, 0x5A, 0x3C)),
]
for label, l, t, color in boxes_flow:
    rect(s, l, t, 3.6, 1.5, color)
    txt(s, label, l + 0.15, t + 0.1, 3.3, 1.3, size=14, bold=True,
        color=C_WHITE, align=PP_ALIGN.CENTER)

txt(s, "→", 3.95, 1.45, 0.5, 0.5, size=28, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s, "→", 8.2, 1.45, 0.5, 0.5, size=28, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(s, "linguagem natural", 0.3, 2.55, 3.6, 0.3,
    size=10, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)
txt(s, "implementação dos passos", 4.5, 2.55, 3.6, 0.3,
    size=10, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)
txt(s, "ações no browser", 8.7, 2.55, 3.6, 0.3,
    size=10, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)

divider(s, 3.0, C_GREEN)

# feature file
txt(s, "alertas.feature", 0.3, 3.15, 6.0, 0.35, size=13, bold=True, color=C_CODE_YLW)
code_block(s, [
    ("# language: pt", C_CODE_CMT),
    ("Funcionalidade: Caixas de Alerta JavaScript", C_CODE_YLW),
    ("  Como um estudante de QA", C_CODE_CMT),
    ("  Quero interagir com os alertas", C_CODE_CMT),
    ("  Para aprender a lidar com diálogos nativos", C_CODE_CMT),
    ("", C_WHITE),
    ("  Contexto:", C_CODE_BLU),
    ("    Dado que acesso a página de Caixas de Alerta", C_WHITE),
    ("", C_WHITE),
    ("  Cenário: Aceitar a confirmação exibe OK", C_CODE_YLW),
    ("    Quando clico na aba de alerta de índice 1", C_WHITE),
    ("    E clico no botão Tentar", C_WHITE),
    ("    E aceito a caixa de confirmação", C_WHITE),
    ("    Então o resultado deve conter \"You pressed OK!\"", C_CODE_GRN),
], 0.3, 3.55, 6.0, 3.65)

# step definitions
txt(s, "alertas.steps.js", 6.8, 3.15, 6.2, 0.35, size=13, bold=True, color=C_ACCENT)
code_block(s, [
    ("import { When, Then } from", C_CODE_YLW),
    ("  '@badeball/cypress-cucumber-preprocessor'", C_CODE_YLW),
    ("import AlertasPage from '../../pages/AlertasPage'", C_CODE_YLW),
    ("", C_WHITE),
    ("When('clico na aba de alerta de índice {int}',", C_CODE_BLU),
    ("  (indice) => {", C_WHITE),
    ("    AlertasPage.clicarAba(indice)", C_CODE_GRN),
    ("  })", C_WHITE),
    ("", C_WHITE),
    ("When('aceito a caixa de confirmação', () => {", C_CODE_BLU),
    ("  cy.on('window:confirm', () => true)", C_CODE_GRN),
    ("})", C_WHITE),
    ("", C_WHITE),
    ("Then('o resultado deve conter {string}', (msg) => {", C_CODE_BLU),
    ("  AlertasPage.verificarResultadoConfirmacao(msg)", C_CODE_GRN),
    ("})", C_WHITE),
], 6.8, 3.55, 6.2, 3.65)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — CYPRESS: COMANDOS E EXECUÇÃO
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_GREEN)
txt(s, "Como Executar — Cypress", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "CYPRESS — RUN", C_GREEN)

# steps
steps = [
    ("1", "Clone o repositório", "git clone <url-do-repositório>"),
    ("2", "Acesse a pasta",      "cd cypress-project"),
    ("3", "Instale as deps",     "npm install"),
]
for i, (num, label, cmd) in enumerate(steps):
    lx = 0.3 + i * 4.3
    rect(s, lx, 0.85, 4.0, 1.3, C_ACCENT)
    txt(s, num, lx + 0.18, 0.9, 0.5, 0.45, size=28, bold=True, color=C_WHITE)
    txt(s, label, lx + 0.7, 0.92, 3.1, 0.4, size=14, bold=True, color=C_LIGHT)
    txt(s, cmd, lx + 0.18, 1.38, 3.7, 0.6, size=12, color=C_CODE_GRN, italic=True)

divider(s, 2.35, C_GREEN)

# commands table
txt(s, "Comandos disponíveis", 0.3, 2.5, 12.0, 0.4, size=16, bold=True, color=C_WHITE)
add_table(s,
    ["Comando npm", "O que faz", "Modo"],
    [
        ["npm test",             "Executa todos os testes Cypress", "Headless"],
        ["npm run test:sem-bdd", "Apenas testes Page Object Model", "Headless"],
        ["npm run test:bdd",     "Apenas testes BDD/Gherkin",       "Headless"],
        ["npm run cypress:abrir","Abre a interface gráfica do Cypress", "Interativo"],
        ["npm run abrir:sem-bdd","Interface gráfica — somente POM", "Interativo"],
        ["npm run abrir:bdd",    "Interface gráfica — somente BDD", "Interativo"],
    ],
    0.3, 2.95, 12.7, 3.25,
    hdr_bg=C_GREEN, row_bg=C_BLUE, row_alt=C_DARK
)

txt(s, "⚙️  Pré-requisito:  Node.js 20+  (verif. com:  node -v)", 0.3, 6.35, 12.7, 0.4,
    size=12, color=C_GRAY, italic=True)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — API TESTING: CONCEITOS
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_ORANGE)
txt(s, "Projeto 2 — Testes de API", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "API TESTING", C_ORANGE)

# http flow diagram
rect(s, 0.3, 0.85, 12.73, 2.1, C_BLUE)
txt(s, "Como funciona uma requisição HTTP", 0.5, 0.9, 8, 0.4,
    size=14, bold=True, color=C_LIGHT)

# boxes in flow
def http_box(slide, l, t, label, sublabel, color):
    rect(slide, l, t, 2.4, 1.25, color)
    txt(slide, label,    l+0.15, t+0.15, 2.1, 0.45, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sublabel, l+0.15, t+0.65, 2.1, 0.5,  size=10, color=C_LIGHT, align=PP_ALIGN.CENTER)

http_box(s, 0.4,  1.3, "Cliente",      "Newman / Postman\nfaz a requisição", C_ACCENT)
txt(s, "GET /pokemon/pikachu", 2.85, 1.7, 2.7, 0.45, size=11, color=C_CODE_YLW, align=PP_ALIGN.CENTER, italic=True)
txt(s, "──────────────►", 2.85, 1.55, 2.7, 0.4, size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
http_box(s, 5.55, 1.3, "Servidor",     "pokeapi.co\nprocessa o pedido", RGBColor(0x3A, 0x1C, 0x71))
txt(s, "◄──────────────", 7.98, 1.55, 2.7, 0.4, size=18, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(s, '200 OK  {"name":"pikachu"}', 7.98, 1.7, 2.7, 0.45, size=11, color=C_CODE_GRN, align=PP_ALIGN.CENTER, italic=True)
http_box(s, 10.7, 1.3, "Resposta",    "JSON com\ndados do Pokémon", RGBColor(0x05, 0x5A, 0x3C))

divider(s, 3.15, C_ORANGE)

# concepts
txt(s, "Conceitos cobertos nos testes", 0.3, 3.3, 9, 0.38, size=15, bold=True, color=C_WHITE)

concepts = [
    ("📡  Status Codes",    "200 OK, 404 Not Found — toda resposta tem um código de status",     C_GREEN),
    ("🔍  Assertions",      "pm.expect(resposta.name).to.equal('pikachu') — valida os dados",     C_ACCENT),
    ("🌐  Variáveis de Env","{{baseUrl}} evita repetição de URL em todos os requests",              C_ORANGE),
    ("📄  JSON Response",   "Resposta estruturada em JSON — validamos campos, tipos e valores",     RGBColor(0x3A, 0x1C, 0x71)),
    ("📃  Paginação",       "?limit=10 — parâmetros de query string para controlar resultados",      RGBColor(0x05, 0x5A, 0x3C)),
    ("⏱️  Performance",     "pm.response.responseTime < 3000 — tempo de resposta aceitável",         RGBColor(0x5A, 0x38, 0x00)),
]
for i, (title, desc, color) in enumerate(concepts):
    col = i % 3
    row = i // 3
    lx = 0.3 + col * 4.3
    ty = 3.75 + row * 1.38
    rect(s, lx, ty, 4.0, 1.22, color)
    txt(s, title, lx+0.15, ty+0.1, 3.7, 0.38, size=13, bold=True, color=C_WHITE)
    txt(s, desc,  lx+0.15, ty+0.5, 3.7, 0.65, size=10, color=C_LIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — POKEAPI: O QUE É TESTADO
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_ORANGE)
txt(s, "PokéAPI — O que é testado?", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "API — TESTES", C_ORANGE)

add_table(s,
    ["#", "Request", "Pokémon", "Validações", "Status"],
    [
        ["1", "GET /pokemon/pikachu",     "Pikachu ⚡",   "status 200, name=='pikachu', id==25, types não vazio, respTime<3s", "✅ 200"],
        ["2", "GET /pokemon/bulbasaur",   "Bulbasaur 🌿",  "status 200, id==1, height>0, weight>0, len(types)==2",             "✅ 200"],
        ["3", "GET /pokemon/charmander",  "Charmander 🔥", "status 200, abilities não vazio, tipo=='fire'",                    "✅ 200"],
        ["4", "GET /pokemon/9999",        "Inexistente ❌", "status 404, NÃO deve ser 200",                                    "✅ 404"],
        ["5", "GET /pokemon?limit=10",    "Listagem 📋",    "status 200, count>0, len(results)==10, next existe",               "✅ 200"],
    ],
    0.3, 0.82, 12.7, 3.45,
    hdr_bg=C_ORANGE
)

divider(s, 4.45, C_ORANGE)

# assertion anatomy
txt(s, "Anatomia de uma assertion (pm.test)", 0.3, 4.6, 12.0, 0.4, size=15, bold=True, color=C_WHITE)
code_block(s, [
    ("pm.test('Nome do Pokémon deve ser pikachu',  function () {", C_CODE_BLU),
    ("  // ① Extrai o corpo da resposta como objeto JavaScript", C_CODE_CMT),
    ("  const resposta = pm.response.json()", C_CODE_GRN),
    ("", C_WHITE),
    ("  // ② Verifica que o campo 'name' é igual a 'pikachu'", C_CODE_CMT),
    ("  pm.expect(resposta.name).to.equal('pikachu')", C_CODE_YLW),
    ("})", C_WHITE),
], 0.3, 5.05, 12.7, 2.2)

# labels
txt(s, "① nome do teste", 0.5, 7.0, 3.0, 0.35, size=10, color=C_GRAY, italic=True)
txt(s, "② extrai o JSON", 4.5, 7.0, 3.0, 0.35, size=10, color=C_GRAY, italic=True)
txt(s, "③ assertion (validação)", 8.5, 7.0, 4.0, 0.35, size=10, color=C_GRAY, italic=True)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — POSTMAN: COLLECTION E ENVIRONMENT
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_ORANGE)
txt(s, "Postman — Collection & Environment", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "POSTMAN", C_ORANGE)

# left: collection structure
rect(s, 0.3, 0.85, 5.8, 6.35, C_BLUE)
txt(s, "📂  pokeapi.collection.json", 0.5, 0.95, 5.4, 0.4, size=13, bold=True, color=C_ORANGE)

collection_tree = [
    ("📁  PokéAPI — Testes de API",         C_WHITE,    True,  13),
    ("  └── 📁  Pokémon",                   C_CODE_YLW, True,  12),
    ("        ├── 📄  Buscar Pikachu",       C_LIGHT,    False, 11),
    ("        │     └── Scripts → pm.test()",C_CODE_CMT, False, 10),
    ("        ├── 📄  Buscar Bulbasaur",     C_LIGHT,    False, 11),
    ("        │     └── Scripts → pm.test()",C_CODE_CMT, False, 10),
    ("        ├── 📄  Buscar Charmander",    C_LIGHT,    False, 11),
    ("        │     └── Scripts → pm.test()",C_CODE_CMT, False, 10),
    ("        ├── 📄  Pokémon Inexistente",  C_LIGHT,    False, 11),
    ("        │     └── Scripts → pm.test()",C_CODE_CMT, False, 10),
    ("        └── 📄  Listar Pokémon",       C_LIGHT,    False, 11),
    ("              └── Scripts → pm.test()",C_CODE_CMT, False, 10),
]
tb2 = s.shapes.add_textbox(Inches(0.45), Inches(1.4), Inches(5.5), Inches(5.6))
tb2.word_wrap = False
tf2 = tb2.text_frame
tf2.word_wrap = False
first = True
for text, color, bold, size in collection_tree:
    if first:
        p = tf2.paragraphs[0]; first = False
    else:
        p = tf2.add_paragraph()
    run = p.add_run()
    run.text = text; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = "Courier New"

# right: environment + package.json
rect(s, 6.4, 0.85, 6.6, 2.6, RGBColor(0x0D, 0x22, 0x42))
txt(s, "🌐  pokeapi.environment.json", 6.6, 0.95, 6.2, 0.38, size=13, bold=True, color=C_ORANGE)
code_block(s, [
    ("{", C_WHITE),
    ('  "name": "PokéAPI — Ambiente de Testes",', C_CODE_YLW),
    ('  "values": [', C_WHITE),
    ('    {', C_WHITE),
    ('      "key": "baseUrl",', C_CODE_BLU),
    ('      "value": "https://pokeapi.co/api/v2"', C_CODE_GRN),
    ('    }', C_WHITE),
    ('  ]', C_WHITE),
    ('}', C_WHITE),
], 6.4, 1.38, 6.6, 2.0)

rect(s, 6.4, 3.65, 6.6, 3.55, RGBColor(0x0D, 0x22, 0x42))
txt(s, "📦  package.json", 6.6, 3.72, 6.2, 0.38, size=13, bold=True, color=C_ORANGE)
code_block(s, [
    ('"scripts": {', C_WHITE),
    ('  "test": "newman run',          C_CODE_BLU),
    ('    collections/pokeapi.collection.json', C_CODE_YLW),
    ('    -e environments/pokeapi.environment.json"', C_CODE_YLW),
    ('},', C_WHITE),
    ('"devDependencies": {', C_WHITE),
    ('  "newman": "^6.1.0",', C_CODE_GRN),
    ('  "newman-reporter-htmlextra": "^1.22.11"', C_CODE_GRN),
    ('}', C_WHITE),
], 6.4, 4.15, 6.6, 2.9)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — API TESTING: EXECUÇÃO
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_ORANGE)
txt(s, "Como Executar — API Testing", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "API — RUN", C_ORANGE)

# steps
steps2 = [
    ("1", "Acesse a pasta",  "cd api-testing"),
    ("2", "Instale as deps", "npm install"),
    ("3", "Execute!",        "npm test"),
]
for i, (num, label, cmd) in enumerate(steps2):
    lx = 0.3 + i * 4.3
    rect(s, lx, 0.85, 4.0, 1.3, C_ORANGE)
    txt(s, num, lx + 0.18, 0.9, 0.5, 0.45, size=28, bold=True, color=C_WHITE)
    txt(s, label, lx + 0.7, 0.92, 3.1, 0.4, size=14, bold=True, color=C_LIGHT)
    txt(s, cmd, lx + 0.18, 1.38, 3.7, 0.6, size=12, color=C_CODE_GRN, italic=True)

divider(s, 2.35, C_ORANGE)

# terminal output simulation
txt(s, "Saída do terminal  (npm test)", 0.3, 2.5, 6.0, 0.38, size=14, bold=True, color=C_WHITE)
code_block(s, [
    ("newman run collections/pokeapi.collection.json ...", C_CODE_CMT),
    ("", C_WHITE),
    ("PokéAPI — Testes de API", C_CODE_YLW),
    ("", C_WHITE),
    ("→ Pokémon", C_WHITE),
    ("  ✓ Buscar Pokémon por nome — Pikachu (312ms)", C_CODE_GRN),
    ("  ✓ Buscar Pokémon por nome — Bulbasaur (289ms)", C_CODE_GRN),
    ("  ✓ Buscar Pokémon por nome — Charmander (301ms)", C_CODE_GRN),
    ("  ✓ Buscar Pokémon inexistente — ID 9999 (198ms)", C_CODE_GRN),
    ("  ✓ Listar Pokémon — Paginação com 10 resultados (276ms)", C_CODE_GRN),
    ("", C_WHITE),
    ("┌─────────────────────────────────────────────┐", C_GRAY),
    ("│          5 requests │ 0 failures             │", C_CODE_GRN),
    ("│         22 assertions │ 0 failed              │", C_CODE_GRN),
    ("└─────────────────────────────────────────────┘", C_GRAY),
], 0.3, 2.95, 6.0, 4.25)

# report
txt(s, "Relatório HTML  (npm run test:relatorio)", 6.6, 2.5, 6.4, 0.38, size=14, bold=True, color=C_WHITE)
code_block(s, [
    ("newman run collections/pokeapi.collection.json", C_CODE_CMT),
    ("  -e environments/pokeapi.environment.json", C_CODE_CMT),
    ("  --reporters cli,htmlextra", C_CODE_YLW),
    ("  --reporter-htmlextra-export relatorio.html", C_CODE_YLW),
], 6.6, 2.95, 6.4, 1.4)

rect(s, 6.6, 4.45, 6.4, 2.75, C_BLUE)
txt(s, "📊  relatorio.html gerado!", 6.8, 4.55, 6.0, 0.4, size=14, bold=True, color=C_ORANGE)
items_report = [
    "• Abre no browser — visual e detalhado",
    "• Lista todos os requests e respostas",
    "• Mostra assertions ✅ / ❌ por teste",
    "• Exibe tempo de resposta de cada request",
    "• Ideal para apresentar resultados ao time",
]
for i, item in enumerate(items_report):
    txt(s, item, 6.8, 5.05 + i * 0.41, 6.0, 0.38, size=12, color=C_LIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — CI/CD: GITHUB ACTIONS
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_ACCENT)
txt(s, "CI/CD — GitHub Actions", 0.4, 0.1, 10, 0.55, size=26, bold=True)
section_badge(s, "CI/CD", C_ACCENT)

# trigger
rect(s, 0.3, 0.85, 12.7, 0.9, RGBColor(0x0D, 0x22, 0x42))
txt(s, "⚡  Disparador:  push para  main  ou  develop  |  Pull Request para  main",
    0.5, 0.95, 12.3, 0.55, size=14, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# three jobs
def job_card(slide, l, t, title, icon, color, steps_list):
    rect(slide, l, t, 3.9, 5.4, color)
    txt(slide, icon, l+0.15, t+0.12, 0.7, 0.55, size=28)
    txt(slide, title, l+0.85, t+0.15, 2.9, 0.55, size=14, bold=True, color=C_WHITE)
    for i, step in enumerate(steps_list):
        ic, label = step
        txt(slide, ic, l+0.15, t+0.85+i*0.72, 0.4, 0.5, size=14)
        txt(slide, label, l+0.55, t+0.88+i*0.72, 3.2, 0.6, size=11, color=C_LIGHT)

job_card(s, 0.3, 1.85, "Cypress\nSem BDD (POM)", "🌲", RGBColor(0x02, 0x5A, 0x30), [
    ("①", "Checkout do repositório"),
    ("②", "Configura Node.js 20"),
    ("③", "cd cypress-project"),
    ("④", "npm ci"),
    ("⑤", "npm run test:sem-bdd"),
    ("📸", "Upload screenshots se falhar"),
])

job_card(s, 4.7, 1.85, "Cypress\nBDD (Gherkin)", "📋", RGBColor(0x3A, 0x1C, 0x71), [
    ("①", "Checkout do repositório"),
    ("②", "Configura Node.js 20"),
    ("③", "cd cypress-project"),
    ("④", "npm ci"),
    ("⑤", "npm run test:bdd"),
    ("📸", "Upload screenshots se falhar"),
])

job_card(s, 9.1, 1.85, "Newman\n(PokéAPI)", "📮", RGBColor(0x7A, 0x2D, 0x00), [
    ("①", "Checkout do repositório"),
    ("②", "Configura Node.js 20"),
    ("③", "cd api-testing"),
    ("④", "npm ci"),
    ("⑤", "npm test"),
    ("", ""),
])

# parallel label
txt(s, "◄─────────────────  3 jobs em paralelo  ─────────────────►",
    0.3, 7.1, 12.7, 0.35, size=12, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — QUADRO COMPARATIVO
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.72, C_DARK)
divider(s, 0, C_ACCENT, 0, 13.33)
txt(s, "Quadro Comparativo — Cypress vs API Testing", 0.4, 0.1, 11, 0.55, size=24, bold=True)

add_table(s,
    ["",             "Cypress E2E",               "API Testing (Newman)"],
    [
        ["Objetivo",     "Simula o usuário no browser",    "Valida respostas HTTP da API"],
        ["Alvo",         "globalsqa.com (frontend)",        "pokeapi.co (API REST)"],
        ["Linguagem",    "JavaScript",                      "JavaScript (scripts Postman)"],
        ["Arquivo base", "*.spec.js / *.feature",          "*.collection.json"],
        ["Execução",     "npm test / npm run test:bdd",    "npm test"],
        ["CI/CD jobs",   "2 jobs (sem-bdd + bdd)",         "1 job (testes-api)"],
        ["Relatório",    "Screenshots em caso de falha",    "HTML via htmlextra"],
        ["Conceito-chave","Page Object Model + BDD",        "Assertions + Status Codes"],
    ],
    0.3, 0.85, 12.7, 5.3,
    hdr_bg=C_ACCENT,
    row_bg=RGBColor(0x10, 0x22, 0x40),
    row_alt=RGBColor(0x07, 0x13, 0x28)
)

txt(s, "Ambos os projetos são escritos em JavaScript e integrados ao mesmo pipeline de CI/CD.",
    0.3, 6.35, 12.7, 0.4, size=13, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — ENCERRAMENTO
# ──────────────────────────────────────────────────────────────────────────────
s = blank_slide()
bg(s, C_DARK)
rect(s, 0, 0, 0.55, 7.5, C_ACCENT)
rect(s, 0.55, 0, 0.12, 7.5, RGBColor(0x12, 0x4A, 0x9A))

txt(s, "Resumo", 1.2, 1.0, 11, 0.9, size=56, bold=True, color=C_WHITE)
divider(s, 2.1, C_ACCENT, 1.2, 10.5)

summary = [
    ("🌲  cypress-project/",    "Testes E2E em 2 estilos (POM e BDD) no site GlobalSQA",       C_GREEN),
    ("📮  api-testing/",        "5 testes de API com Postman/Newman na PokéAPI",                C_ORANGE),
    ("⚡  .github/workflows/",   "Pipeline CI/CD automático com 3 jobs em paralelo",             C_ACCENT),
    ("📄  README.md",           "Documentação e guia de início rápido para os dois projetos",    C_GRAY),
]
for i, (title, desc, color) in enumerate(summary):
    rect(s, 1.2, 2.35 + i * 1.15, 11.0, 1.0, C_BLUE)
    rect(s, 1.2, 2.35 + i * 1.15, 0.18, 1.0, color)
    txt(s, title, 1.55, 2.42 + i * 1.15, 4.5, 0.42, size=14, bold=True, color=C_WHITE)
    txt(s, desc,  1.55, 2.82 + i * 1.15, 10.4, 0.42, size=12, color=C_LIGHT)

txt(s, "S07 — Qualidade de Software  |  QA DevOps",
    1.2, 7.0, 11, 0.38, size=12, color=C_GRAY, italic=True)


# ──────────────────────────────────────────────────────────────────────────────
# SALVAR
# ──────────────────────────────────────────────────────────────────────────────
OUTPUT = "/home/user/S07/S07_Qualidade_de_Software.pptx"
prs.save(OUTPUT)
print(f"✅  Apresentação salva em: {OUTPUT}")
print(f"    Slides: {len(prs.slides)}")
